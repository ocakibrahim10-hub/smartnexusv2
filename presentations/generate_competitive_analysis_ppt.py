"""SmartNexus Türkiye ERP Rekabet Analizi — PowerPoint oluşturucu."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# SmartNexus Hosting tema renkleri
PRIMARY = RGBColor(96, 107, 223)       # #606BDF
PRIMARY_DARK = RGBColor(72, 82, 190)
BG_LIGHT = RGBColor(251, 248, 255)    # #FBF8FF
TEXT_DARK = RGBColor(27, 27, 31)      # #1B1B1F
TEXT_MUTED = RGBColor(70, 70, 79)     # #46464F
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(239, 237, 244)      # #EFEDF4
SUCCESS = RGBColor(31, 138, 101)
WARNING = RGBColor(192, 133, 50)
DANGER = RGBColor(207, 45, 86)
ACCENT_LIGHT = RGBColor(224, 224, 255)  # #E0E0FF

OUT_PATH = Path(__file__).resolve().parent / "SmartNexus_Rekabet_Analizi_2026.pptx"


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = ACCENT_LIGHT
        p2.space_before = Pt(4)


def add_footer(slide, text: str = "SmartNexus · Gizli · Haziran 2026"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullet_box(slide, left, top, width, height, items, font_size=14, color=TEXT_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return box


def add_kpi_card(slide, left, top, width, height, label, value, accent=False):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    if accent:
        accent_bar = slide.shapes.add_shape(1, left, top, width, Inches(0.06))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = PRIMARY
        accent_bar.line.fill.background()
    tf = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.3)).text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY if accent else TEXT_DARK
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(4)


def add_table(slide, left, top, width, rows, col_widths=None):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, Inches(0.38 * len(rows)))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if r > 0 else 11)
                p.font.bold = r == 0
                if r == 0:
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = TEXT_DARK
                    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_DARK
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- 1 Kapak ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    band = s.shapes.add_shape(1, Inches(0), Inches(2.4), Inches(13.333), Inches(2.6))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()
    title = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "SmartNexus Rekabet Analizi"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.8), Inches(3.85), Inches(11), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Türkiye ERP & Bulut Muhasebe Pazarı · Haziran 2026"
    sp.font.size = Pt(18)
    sp.font.color.rgb = ACCENT_LIGHT
    meta = s.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8), Inches(1))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "SmartERP / SmartNexus Platform\nKapsamlı özellik · fiyat · entegrasyon karşılaştırması"
    mp.font.size = Pt(14)
    mp.font.color.rgb = TEXT_MUTED

    # --- 2 Yönetici özeti ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Yönetici Özeti")
    add_kpi_card(s, Inches(0.6), Inches(1.35), Inches(2.9), Inches(1.15), "Türkiye geneli sıra", "~21 / 35", accent=True)
    add_kpi_card(s, Inches(3.7), Inches(1.35), Inches(2.9), Inches(1.15), "Ağırlıklı puan", "52 / 100")
    add_kpi_card(s, Inches(6.8), Inches(1.35), Inches(2.9), Inches(1.15), "Niş segment sırası", "3–5 / 8")
    add_kpi_card(s, Inches(9.9), Inches(1.35), Inches(2.8), Inches(1.15), "Ürün tamamlanma", "~%58")
    add_bullet_box(s, Inches(0.6), Inches(2.75), Inches(12), Inches(3.8), [
        "SmartNexus: muhasebe + stok + POS + TMS + B2B + pazaryeri + çoklu şube + bayi kanalı tek SaaS çatısında.",
        "Fiyat bandı: ₺499 – ₺2.999 / ay — geleneksel ERP’ye göre %70–90 daha düşük TCO.",
        "Güçlü yönler: modern UX (%88), bayi mimarisi (%75), modüler paketleme, zengin demo verisi.",
        "Kritik açıklar: gerçek GİB e-fatura (%18), pazaryeri mock (%20), genel muhasebe yok, canlı prod eksik.",
        "12 ay hedef: sıra 12–15, puan 68–72, 50–100 ödeme yapan müşteri, ₺75K–150K MRR.",
    ], font_size=13)
    add_footer(s)

    # --- 3 Metodoloji ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Değerlendirme Metodolojası", "10 kriter · 0–100 puan · Türkiye pazarı ağırlıklı")
    criteria = [
        ("Muhasebe & finans", "%15"),
        ("Stok & WMS", "%12"),
        ("Satış & POS", "%10"),
        ("Lojistik (TMS)", "%8"),
        ("B2B & e-ticaret", "%10"),
        ("E-dönüşüm (GİB)", "%12"),
        ("Çoklu şube / bayi", "%10"),
        ("Fiyat / performans", "%10"),
        ("UX & modernlik", "%8"),
        ("Pazar olgunluğu", "%15"),
    ]
    rows = [["Kriter", "Ağırlık"]] + [[a, b] for a, b in criteria]
    add_table(s, Inches(0.6), Inches(1.4), Inches(5.5), rows, [Inches(3.8), Inches(1.5)])
    add_bullet_box(s, Inches(6.5), Inches(1.4), Inches(6.2), Inches(5.2), [
        "Veri kaynakları:",
        "• SmartNexus kod tabanı (schema, API, UI, seed)",
        "• Rakip resmi fiyat listeleri (Logo, Mikro, Paraşüt, KolayBi, Bizim Hesap)",
        "• Sektör inceleme raporları (2025–2026)",
        "",
        "Karşılaştırılan rakipler:",
        "Logo Tiger · Nebim V3 · Mikro Jump · Dia · Paraşüt",
        "KolayBi · Bizim Hesap · Logo İşbaşı · Uyumsoft",
    ], font_size=12)
    add_footer(s)

    # --- 4 Ürün envanteri ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "SmartNexus Ürün Envanteri", "8 modül grubu · 25 alt modül")
    pkg_rows = [
        ["Paket", "Aylık", "Kapsam"],
        ["Basic", "₺499", "Fatura, cari, kasa, ürün, depo, POS"],
        ["Professional", "₺1.499", "Basic + e-dönüşüm, stok, TMS, B2B, pazaryeri"],
        ["Platinum", "₺2.999", "Pro + AI tahmin, çoklu şube, şube takibi"],
        ["Bayi Platform", "₺2.999", "İşletme, komisyon, faturalama, abonelik"],
    ]
    add_table(s, Inches(0.6), Inches(1.35), Inches(7.2), pkg_rows, [Inches(1.8), Inches(1.2), Inches(4.0)])
    mod_rows = [
        ["Modül", "Tamamlanma"],
        ["Cari & Fatura", "%82"],
        ["Kasa & Banka", "%81"],
        ["Stok & Depo", "%80"],
        ["POS", "%73"],
        ["B2B", "%65"],
        ["TMS", "%58"],
        ["E-Dönüşüm", "%18"],
        ["Pazaryeri", "%20"],
    ]
    add_table(s, Inches(8.1), Inches(1.35), Inches(4.6), mod_rows, [Inches(2.8), Inches(1.6)])
    add_footer(s)

    # --- 5 Pazar haritası ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Türkiye ERP Pazar Haritası")
    tiers = [
        ("Kurumsal ERP", "₺100K – ₺1M+ lisans", "Logo Tiger · Nebim V3 · SAP B1 · Dynamics 365"),
        ("KOBİ Ticari ERP", "₺15K – ₺150K / yıl", "Mikro Jump · Logo Netsis · Dia · ETA"),
        ("Bulut Ön Muhasebe", "₺300 – ₺1.200 / ay", "Paraşüt · KolayBi · Bizim Hesap · İşbaşı"),
        ("Entegre Modüler SaaS", "₺500 – ₺5.000 / ay", "SmartNexus · Dia Bulut · Logo kanal"),
    ]
    y = 1.45
    for title, price, players in tiers:
        card = s.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.1), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        tf = s.shapes.add_textbox(Inches(0.85), Inches(y + 0.12), Inches(11.5), Inches(0.9)).text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = f"{title}  ·  {price}"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = players
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)
        y += 1.35
    highlight = s.shapes.add_shape(1, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.06))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = PRIMARY
    highlight.line.fill.background()
    add_footer(s)

    # --- 6 Özellik matrisi ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Özellik Karşılaştırması", "Pazar lideri seviyesine göre yetkinlik %")
    feat_rows = [
        ["Alan", "SmartNexus", "Logo", "Mikro", "Paraşüt", "Nebim"],
        ["Cari & Fatura", "%82", "%95", "%90", "%90", "%88"],
        ["Stok & Depo", "%80", "%90", "%88", "%45", "%95"],
        ["POS", "%73", "%85", "%80", "%20", "%98"],
        ["E-Fatura (GİB)", "%18", "%95", "%90", "%95", "%88"],
        ["Pazaryeri", "%20", "%60", "%50", "%55", "%90"],
        ["Çoklu Şube", "%75", "%90", "%75", "%10", "%95"],
        ["Bayi SaaS", "%75", "%40", "%35", "%5", "%50"],
        ["UX Modernlik", "%88", "%55", "%60", "%85", "%65"],
        ["Ağırlıklı Ort.", "%52", "%83", "%78", "%58", "%84"],
    ]
    add_table(s, Inches(0.5), Inches(1.3), Inches(12.3), feat_rows,
              [Inches(2.4), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)])
    add_footer(s)

    # --- 7 Fiyatlandırma ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Fiyatlandırma Karşılaştırması", "KOBİ · 5 kullanıcı · temel modüller")
    price_rows = [
        ["Ürün", "Başlangıç / ay", "Orta / ay", "Üst / ay", "Model"],
        ["SmartNexus", "₺499", "₺1.499", "₺2.999", "Abonelik"],
        ["Paraşüt", "~₺752+KDV", "~₺940+KDV", "—", "Yıllık paket"],
        ["KolayBi", "~₺570", "—", "—", "Abonelik"],
        ["Bizim Hesap", "~₺719+KDV", "~₺920+KDV", "—", "Yıllık"],
        ["Logo İşbaşı", "~₺342", "—", "—", "Bulut"],
        ["Mikro Jump*", "~₺4.400", "—", "—", "Lisans/yıl"],
        ["Logo Tiger*", "~₺29.000", "—", "—", "Lisans/yıl"],
    ]
    add_table(s, Inches(0.5), Inches(1.35), Inches(12.3), price_rows,
              [Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)])
    note = s.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12), Inches(0.8))
    np = note.text_frame.paragraphs[0]
    np.text = "* Yıllık lisansın 12 aya bölünmüş yaklaşık değeri · kurulum/danışmanlık hariç"
    np.font.size = Pt(10)
    np.font.color.rgb = TEXT_MUTED
    add_kpi_card(s, Inches(0.6), Inches(5.5), Inches(3.8), Inches(1.2), "SmartNexus vs Logo Tiger TCO", "%70–90 ucuz", accent=True)
    add_kpi_card(s, Inches(4.6), Inches(5.5), Inches(3.8), Inches(1.2), "5 yıl SmartNexus Platinum", "~₺180K")
    add_kpi_card(s, Inches(8.6), Inches(5.5), Inches(4.1), Inches(1.2), "5 yıl Logo Tiger TCO", "₺2M+")
    add_footer(s)

    # --- 8 TCO ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "5 Yıllık Toplam Sahip Olma Maliyeti (TCO)")
    tco_rows = [
        ["Kalem", "SmartNexus", "Paraşüt", "Mikro Jump", "Logo Tiger"],
        ["Yazılım (5 yıl)", "₺180.000", "₺90K–150K", "₺260K+", "₺1.75M+"],
        ["Kurulum", "₺0–20K", "₺0–5K", "₺50K–150K", "₺200K–500K"],
        ["E-kontör / ek", "Dahil*", "₺20K–60K", "₺30K", "₺80K+"],
        ["Altyapı", "Dahil", "Dahil", "₺30K", "₺50K+"],
        ["TOPLAM", "₺180K–200K", "₺110K–215K", "₺370K–500K", "₺2M+"],
    ]
    add_table(s, Inches(0.6), Inches(1.4), Inches(12), tco_rows,
              [Inches(2.5), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)])
    add_bullet_box(s, Inches(0.6), Inches(4.2), Inches(12), Inches(2.5), [
        "SmartNexus: entegre paket — stok, POS, TMS, B2B tek abonelikte; geleneksel ERP modül lisansı gerektirmez.",
        "Paraşüt: dar kapsam (ön muhasebe) ama olgun e-fatura ve banka entegrasyonu.",
        "Logo/Mikro: yüksek ilk yatırım + danışmanlık; uzun vadede kurumsal ölçekte tercih edilir.",
    ], font_size=12)
    add_footer(s)

    # --- 9 Entegrasyonlar ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Entegrasyon & Ekosistem Boşlukları")
    int_rows = [
        ["Entegrasyon", "SmartNexus", "Pazar", "Gap"],
        ["GİB E-Fatura", "Mock", "Zorunlu", "-%95"],
        ["E-Arşiv / İrsaliye", "Yok", "Yüksek", "-%100"],
        ["Pazaryeri API", "Mock UI", "E-ticaret kritik", "-%100"],
        ["Banka API", "Yok", "Bulut standart", "-%100"],
        ["Ödeme Gateway", "Yok", "B2B gerekli", "-%100"],
        ["Luca / Mali müşavir", "Yok", "TR satış hızlandırıcı", "-%100"],
        ["Hugin Yazar Kasa", "Kısmi", "Perakende", "-%40"],
        ["Mobil Uygulama", "Yok", "Saha satış", "-%100"],
    ]
    add_table(s, Inches(0.6), Inches(1.35), Inches(12), int_rows,
              [Inches(3.2), Inches(2.5), Inches(3.5), Inches(1.5)])
    add_footer(s)

    # --- 10 Puanlama ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "Türkiye Geneli Puanlama", "SmartNexus kriter bazlı analiz")
    score_rows = [
        ["Kriter", "SmartNexus", "Sektör Ort.", "Fark"],
        ["Muhasebe derinliği", "35", "78", "-43"],
        ["Stok & WMS", "80", "82", "-2"],
        ["Satış & POS", "73", "75", "-2"],
        ["TMS", "58", "65", "-7"],
        ["B2B & e-ticaret", "65", "70", "-5"],
        ["E-dönüşüm", "18", "88", "-70"],
        ["Çoklu şube / bayi", "75", "55", "+20"],
        ["Fiyat / performans", "85", "60", "+25"],
        ["UX & modernlik", "88", "68", "+20"],
        ["Pazar olgunluğu", "12", "85", "-73"],
        ["AĞIRLIKLI TOPLAM", "52", "74", "-22"],
    ]
    add_table(s, Inches(0.6), Inches(1.35), Inches(7.5), score_rows,
              [Inches(3.2), Inches(1.4), Inches(1.4), Inches(1.2)])
    rank_rows = [
        ["Sıra", "Ürün", "Puan"],
        ["1", "Logo Tiger 3", "91"],
        ["2", "Nebim V3", "88"],
        ["3", "Paraşüt", "86"],
        ["4", "Mikro Jump", "84"],
        ["5", "Dia Yazılım", "81"],
        ["…", "…", "…"],
        ["~21", "SmartNexus", "52"],
    ]
    add_table(s, Inches(8.4), Inches(1.35), Inches(4.3), rank_rows,
              [Inches(0.7), Inches(2.3), Inches(1.0)])
    add_footer(s)

    # --- 11 SWOT ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "SWOT Analizi")
    swot_data = [
        ("Güçlü Yönler (S)", SUCCESS, [
            "Tek çatıda ERP+WMS+POS+TMS+B2B",
            "Modern SaaS UX (%88)",
            "4 katman bayi mimarisi",
            "Rekabetçi aylık fiyat (₺499+)",
            "Modüler paketleme & zengin demo",
        ]),
        ("Zayıf Yönler (W)", DANGER, [
            "Gerçek GİB e-fatura yok",
            "Genel muhasebe / beyanname yok",
            "Pazaryeri & bayi faturalama mock",
            "Canlı prod & sıfır müşteri",
            "Mobil uygulama yok",
        ]),
        ("Fırsatlar (O)", PRIMARY, [
            "E-dönüşüm zorunluluğu artıyor",
            "KOBİ buluta geçiş trendi",
            "Bayi SaaS marj modeli",
            "Perakende+lojistik entegrasyon talebi",
        ]),
        ("Tehditler (T)", WARNING, [
            "Paraşüt/Dia fiyat baskısı",
            "Logo/Mikro kanal ağı",
            "GİB mevzuat değişiklikleri",
            "Uyumsoft e-fatura+ERP birleşimi",
        ]),
    ]
    positions = [(0.6, 1.35), (6.75, 1.35), (0.6, 4.15), (6.75, 4.15)]
    for (left, top), (title, color, items) in zip(positions, swot_data):
        card = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(5.9), Inches(2.55))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        hdr = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(5.9), Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        ht = s.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.08), Inches(5.5), Inches(0.35))
        ht.text_frame.paragraphs[0].text = title
        ht.text_frame.paragraphs[0].font.size = Pt(13)
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.color.rgb = WHITE
        add_bullet_box(s, Inches(left + 0.2), Inches(top + 0.55), Inches(5.5), Inches(1.9), items, font_size=11)
    add_footer(s)

    # --- 12 Yol haritası ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "İlerleme Yol Haritası", "Öncelik sırasıyla")
    phases = [
        ("Faz 1 · 0–3 Ay", "Canlıya çıkış blokörleri", [
            "GİB e-fatura entegratörü (Uyumsoft/Foriba/Logo)",
            "Production API + DB + Vercel env",
            "Plan limitlerini koda bağla",
            "Rol bazlı yetki (RolesGuard)",
        ]),
        ("Faz 2 · 3–6 Ay", "Rekabetçi eşitlik", [
            "Banka entegrasyonu (5+ banka)",
            "Trendyol + Hepsiburada API",
            "B2B müşteri self-servis portalı",
            "Mobil PWA / React Native",
            "Luca / Excel export",
        ]),
        ("Faz 3 · 6–12 Ay", "Farklılaşma", [
            "Gerçek AI talep tahmini",
            "Ödeme gateway (iyzico/PayTR)",
            "Audit trail & KVKK",
            "Tenant API keys (Platinum)",
            "e-İrsaliye + e-Arşiv",
        ]),
        ("Faz 4 · Sürekli", "Pazar & kanal", [
            "İlk 10 ödeme yapan pilot müşteri",
            "2–3 bayi ortağı (kanal marjı)",
            "Sektör paketleri (Perakende, Lojistik)",
            "Fiyat: Basic agresif + e-kontör geliri",
        ]),
    ]
    x_positions = [0.55, 3.45, 6.35, 9.25]
    for x, (phase, subtitle, items) in zip(x_positions, phases):
        card = s.shapes.add_shape(1, Inches(x), Inches(1.35), Inches(2.75), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        bar = s.shapes.add_shape(1, Inches(x), Inches(1.35), Inches(2.75), Inches(0.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()
        tf = s.shapes.add_textbox(Inches(x + 0.1), Inches(1.42), Inches(2.55), Inches(0.7)).text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = phase
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(9)
        p2.font.color.rgb = ACCENT_LIGHT
        add_bullet_box(s, Inches(x + 0.12), Inches(2.0), Inches(2.5), Inches(4.5), items, font_size=9)
    add_footer(s)

    # --- 13 Hedef metrikler ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, BG_LIGHT)
    add_header_bar(s, "12 Ay Hedef Metrikleri")
    targets = [
        ("Türkiye sırası", "21", "12–15"),
        ("Ağırlıklı puan", "52", "68–72"),
        ("E-dönüşüm puanı", "18", "80+"),
        ("Canlı müşteri", "0", "50–100"),
        ("Aylık MRR", "₺0", "₺75K–150K"),
    ]
    for i, (label, today, target) in enumerate(targets):
        left = Inches(0.6 + (i % 3) * 4.15)
        top = Inches(1.5 + (i // 3) * 2.2)
        add_kpi_card(s, left, top, Inches(3.8), Inches(1.0), f"Bugün: {today}", label, accent=i == 0)
        arrow = s.shapes.add_textbox(left + Inches(1.5), top + Inches(1.05), Inches(0.8), Inches(0.4))
        arrow.text_frame.paragraphs[0].text = "↓"
        arrow.text_frame.paragraphs[0].font.size = Pt(18)
        arrow.text_frame.paragraphs[0].font.color.rgb = PRIMARY
        arrow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tgt = s.shapes.add_shape(1, left, top + Inches(1.35), Inches(3.8), Inches(0.65))
        tgt.fill.solid()
        tgt.fill.fore_color.rgb = ACCENT_LIGHT
        tgt.line.color.rgb = PRIMARY
        tfp = s.shapes.add_textbox(left + Inches(0.15), top + Inches(1.45), Inches(3.5), Inches(0.5)).text_frame
        tfp.paragraphs[0].text = f"Hedef: {target}"
        tfp.paragraphs[0].font.size = Pt(14)
        tfp.paragraphs[0].font.bold = True
        tfp.paragraphs[0].font.color.rgb = PRIMARY_DARK
    conclusion = s.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(1))
    cp = conclusion.text_frame.paragraphs[0]
    cp.text = (
        "Sonuç: SmartNexus vizyon ve mimari olarak güçlü; ticari olgunluk için GİB entegrasyonu + "
        "canlı prod + ilk 10 müşteri kritik öncelik."
    )
    cp.font.size = Pt(14)
    cp.font.color.rgb = TEXT_DARK
    cp.font.bold = True
    add_footer(s)

    # --- 14 Kapanış ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, PRIMARY)
    thanks = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
    tp = thanks.text_frame.paragraphs[0]
    tp.text = "Teşekkürler"
    tp.font.size = Pt(44)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "SmartNexus · Türkiye ERP Rekabet Analizi 2026\nsmartnexus-erp.vercel.app"
    sp.font.size = Pt(16)
    sp.font.color.rgb = ACCENT_LIGHT
    sp.alignment = PP_ALIGN.CENTER

    return prs


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUT_PATH))
    print(f"Olusturuldu: {OUT_PATH}")
    print(f"Slayt sayisi: {len(prs.slides)}")


if __name__ == "__main__":
    main()
